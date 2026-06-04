# -*- coding: utf-8 -*-
"""
Build an editable PowerPoint version of the UnClick investor deck.
Same 19-slide narrative and brand as unclick-investor-deck.html, rebuilt as
native, editable shapes so it opens cleanly in PowerPoint, Keynote, or Slides.

Run:  python3 build_pptx.py
Out:  unclick-investor-deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- brand tokens
TEAL      = RGBColor(0x61, 0xC1, 0xC4)
TEAL_BR   = RGBColor(0x86, 0xDA, 0xDD)
TEAL_DEEP = RGBColor(0x2F, 0x94, 0x99)
NAVY      = RGBColor(0x06, 0x20, 0x2C)
NAVY2     = RGBColor(0x0A, 0x2C, 0x3C)
NAVY3     = RGBColor(0x0E, 0x3A, 0x4E)
INK       = RGBColor(0x0D, 0x2A, 0x36)
BODY      = RGBColor(0x3C, 0x57, 0x63)
MUTED     = RGBColor(0x6F, 0x8B, 0x95)
LIGHTBG   = RGBColor(0xEE, 0xF5, 0xF6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LINE      = RGBColor(0xD6, 0xE4, 0xE8)
RED       = RGBColor(0xD0, 0x53, 0x4A)
GREEN     = RGBColor(0x3F, 0xA9, 0x6A)
AMBER     = RGBColor(0xC8, 0x8F, 0x2A)
DARKTXT   = RGBColor(0xEA, 0xF6, 0xF7)
DARKBODY  = RGBColor(0xBC, 0xD9, 0xDD)
CARD_DARK = RGBColor(0x10, 0x33, 0x44)

FONT = "Segoe UI"     # Windows-native, very close to the brand Inter
MONO = "Consolas"

W, H = 13.333, 7.5
ML   = 0.92           # left margin
CW   = 11.5           # content width

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def _no_shadow(shape):
    # strip python-pptx default preset shadow
    sp = shape._element.spPr
    existing = sp.find(qn('a:effectLst'))
    if existing is None:
        sp.append(sp.makeelement(qn('a:effectLst'), {}))

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=True, radius=0.09):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    _no_shadow(shp)
    return shp

def fullbg(s, color):
    rect(s, 0, 0, W, H, fill=color, rounded=False)

def text(s, l, t, w, h, segments, size, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, name=FONT, spacing=1.06, italic=False, wrap=True):
    """segments: str  OR  list of (txt, color) OR list of (txt, color, bold)."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = align
    try: p.line_spacing = spacing
    except Exception: pass
    if isinstance(segments, str):
        segments = [(segments, None)]
    for seg in segments:
        txt = seg[0]; col = seg[1] if len(seg) > 1 else None
        b = seg[2] if len(seg) > 2 else bold
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = b; r.font.name = name; r.font.italic = italic
        if col is not None: r.font.color.rgb = col
    return tb

def paras(s, l, t, w, h, lines, size, color=BODY, bold=False, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, name=FONT, spacing=1.18, gap=6):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, 0)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        try: p.line_spacing = spacing
        except Exception: pass
        p.space_after = Pt(gap)
        segs = ln if isinstance(ln, list) else [(ln, color, bold)]
        for seg in segs:
            txt = seg[0]; col = seg[1] if len(seg) > 1 else color
            b = seg[2] if len(seg) > 2 else bold
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = b; r.font.name = name
            r.font.color.rgb = col
    return tb

def eyebrow(s, label, y=0.62, dark=False):
    col = TEAL_BR if dark else TEAL_DEEP
    rect(s, ML, y + 0.07, 0.28, 0.035, fill=col, rounded=False)
    text(s, ML + 0.38, y - 0.05, 9, 0.32, [(label.upper(), col)], 11.5, bold=True)

def footer(s, n, dark=False):
    col = RGBColor(0x7F, 0xA6, 0xB0) if not dark else RGBColor(0x6F, 0x97, 0xA1)
    text(s, ML, H - 0.46, 8, 0.3, [("UnClick  ·  Where AI belongs. Humans welcome.", col)], 9.5)
    text(s, W - 2.4, H - 0.46, 1.8, 0.3, [("%02d / 19" % n, col)], 9.5,
         align=PP_ALIGN.RIGHT, name=MONO)

def accent_top(s):
    rect(s, 0, 0, 3.4, 0.055, fill=TEAL, rounded=False)

def logo_wordmark(s, l, t, dark=True, size=26):
    # simple brand lockup: teal cursor glyph block + wordmark
    g = rect(s, l, t, 0.42, 0.42, fill=TEAL, rounded=True, radius=0.25)
    text(s, l, t - 0.02, 0.42, 0.46, [("›", NAVY)], 20, bold=True, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, l + 0.56, t - 0.06, 5, 0.6, [("UnClick", WHITE if dark else INK)], size, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)

def chip(s, l, t, w, h, label, dark=False):
    fill = CARD_DARK if dark else WHITE
    ln = RGBColor(0x1C, 0x4A, 0x59) if dark else LINE
    return rect(s, l, t, w, h, fill=fill, line=ln, line_w=1.0)

def pill(s, l, t, w, h, label, dark=False):
    fill = RGBColor(0x10, 0x33, 0x44) if dark else RGBColor(0xE4, 0xF5, 0xF5)
    bd   = TEAL if dark else RGBColor(0xB9, 0xE6, 0xE7)
    col  = TEAL_BR if dark else TEAL_DEEP
    rect(s, l, t, w, h, fill=fill, line=bd, line_w=1.0, radius=0.5)
    text(s, l, t, w, h, [(label, col)], 11.5, bold=True, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


# ================================================================ SLIDE 1 cover
s = slide(); fullbg(s, NAVY)
rect(s, 6.7, -1.2, 8.5, 7.0, fill=NAVY2, rounded=True, radius=0.5)
rect(s, 8.2, 0.2, 6.0, 5.2, fill=NAVY3, rounded=True, radius=0.5)
logo_wordmark(s, ML, 0.8, dark=True, size=30)
eyebrow(s, "Investor overview  ·  2026", y=1.95, dark=True)
text(s, ML, 2.35, 11.4, 2.0,
     [("The universal remote ", WHITE), ("for AI.", TEAL_BR)], 58, bold=True, spacing=1.0)
text(s, ML, 4.15, 11.0, 0.6, [("Where AI belongs. Humans welcome.", TEAL_BR)], 22, bold=True)
paras(s, ML, 4.85, 9.7, 1.2,
      ["One install gives any AI agent every tool it needs, memory that remembers across "
       "every session and device, and built-in proof the work was done right."],
      16.5, color=DARKBODY, spacing=1.3)
pill(s, ML, 6.15, 5.1, 0.42, "unclick.world  ·  Built in Melbourne, Australia", dark=True)
footer(s, 1, dark=True)

# ================================================================ SLIDE 2 big idea
s = slide(); fullbg(s, NAVY)
rect(s, 7.6, -1.5, 8.0, 7.5, fill=NAVY2, rounded=True, radius=0.5)
eyebrow(s, "The big idea", y=1.15, dark=True)
text(s, ML, 1.7, 11.4, 2.4,
     [("Let your AI stop ", WHITE), ("pretending to be human.", TEAL_BR)], 52, bold=True, spacing=1.0)
paras(s, ML, 4.0, 10.2, 1.6,
      ["Today, AI agents fake their way through screens built for people. They click buttons, "
       "fill in forms, and forget everything the moment you close the tab. UnClick gives them "
       "real rails instead: tools to act, a memory that lasts, and a control tower that proves "
       "every job was done right."],
      17, color=DARKBODY, spacing=1.32)
text(s, ML, 6.15, 11, 0.5, [("Every tool. One install. Every AI app.", WHITE)], 18, bold=True)
footer(s, 2, dark=True)

# ================================================================ SLIDE 3 problem
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "The problem")
text(s, ML, 1.05, 11.4, 1.0, [("AI agents are powerful, but ", INK), ("homeless.", TEAL_DEEP)], 42, bold=True)
paras(s, ML, 2.0, 10.8, 1.0,
      ["Every AI tool reinvents the same missing layer: identity, memory, connections, a "
       "trustworthy toolset, and proof. Switch from Claude to ChatGPT and you start over from zero."],
      16.5, color=BODY, spacing=1.3)
cy, cw, ch = 3.25, 5.55, 2.9
for i,(t,b) in enumerate([
    ("The amnesia tax",
     "Every session starts from scratch. You re-explain your business, your preferences, your "
     "rules. Every time. Your agent made 200 decisions last week and remembers exactly none of them."),
    ("The wrong tool",
     "Asking your AI to work through a website built for people is like doing surgery with oven "
     "mitts. It can sort of work. It is just not the right tool for the job.")]):
    cx = ML + i*(cw+0.4)
    rect(s, cx, cy, cw, ch, fill=WHITE, line=LINE, line_w=1.0)
    text(s, cx+0.35, cy+0.32, cw-0.7, 0.5, [(t, INK)], 20, bold=True)
    paras(s, cx+0.35, cy+1.0, cw-0.7, ch-1.2, [b], 14.5, color=BODY, spacing=1.28)
footer(s, 3)

# ================================================================ SLIDE 4 old vs new
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Why it matters")
text(s, ML, 1.05, 11.4, 1.0, [("The old way is slow, costly, and blind.", INK)], 40, bold=True)
# old card
ox, ow, oy, oh = ML, 4.55, 2.45, 3.0
rect(s, ox, oy, ow, oh, fill=RGBColor(0xFD,0xF1,0xF0), line=RGBColor(0xEA,0xC4,0xC1), line_w=1.2)
text(s, ox+0.32, oy+0.28, ow-0.6, 0.3, [("THE OLD WAY", RED)], 11.5, bold=True)
paras(s, ox+0.32, oy+0.78, ow-0.6, 1.2,
      ["Open browser → Navigate", "→ Fill the form → Click submit", "→ Wait → Read the screen"],
      12.5, color=INK, name=MONO, spacing=1.25, gap=2)
text(s, ox+0.32, oy+2.15, 2.2, 0.5, [("~4,200ms", RED)], 22, bold=True, name=MONO)
text(s, ox+0.32, oy+2.6, 2.4, 0.3, [("PER ACTION", MUTED)], 9.5, bold=True)
text(s, ox+2.4, oy+2.15, 2.0, 0.5, [("~18,000", RED)], 22, bold=True, name=MONO)
text(s, ox+2.4, oy+2.6, 2.0, 0.3, [("TOKENS", MUTED)], 9.5, bold=True)
# badge
text(s, 5.7, 2.9, 2.0, 0.9, [("~110x", TEAL_DEEP)], 40, bold=True, align=PP_ALIGN.CENTER)
paras(s, 5.6, 3.75, 2.2, 0.8, [[("faster", MUTED, True)], [("~150x fewer tokens", MUTED)]],
      10.5, align=PP_ALIGN.CENTER, spacing=1.1, gap=1)
# new card
nx = 7.9; nw = 4.55
rect(s, nx, oy, nw, oh, fill=RGBColor(0xEA,0xFA,0xFB), line=RGBColor(0x9F,0xDE,0xDF), line_w=1.2)
text(s, nx+0.32, oy+0.28, nw-0.6, 0.3, [("THE UNCLICK WAY", TEAL_DEEP)], 11.5, bold=True)
paras(s, nx+0.32, oy+0.85, nw-0.6, 1.0,
      ["POST /v1/schedule/events", [("→ ", INK),("201 Created", INK, True)]],
      13, color=INK, name=MONO, spacing=1.3, gap=3)
text(s, nx+0.32, oy+2.15, 2.0, 0.5, [("~38ms", TEAL_DEEP)], 22, bold=True, name=MONO)
text(s, nx+0.32, oy+2.6, 2.0, 0.3, [("PER ACTION", MUTED)], 9.5, bold=True)
text(s, nx+2.1, oy+2.15, 2.0, 0.5, [("~120", TEAL_DEEP)], 22, bold=True, name=MONO)
text(s, nx+2.1, oy+2.6, 2.0, 0.3, [("TOKENS", MUTED)], 9.5, bold=True)
paras(s, ML, 5.95, 11.4, 0.8,
      ["Same job. One path is built for humans. The other is built for agents. Across thousands "
       "of actions, that gap is the difference between an agent that is a demo and one that is a business."],
      12.5, color=MUTED, spacing=1.25)
footer(s, 4)

# ================================================================ SLIDE 5 why now
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Why now")
text(s, ML, 1.05, 11.4, 1.0, [("The agent wave is here. The rails are not, ", INK), ("yet.", TEAL_DEEP)], 40, bold=True)
cards = [
    ("Agents are exploding",
     "The world is shipping AI agents far faster than the infrastructure that makes them "
     "actually useful. The gap is the opportunity."),
    ("MCP became the standard",
     "Model Context Protocol is the USB-C for AI, now native in Claude, ChatGPT, Cursor and "
     "Copilot. One plug fits every agent, and UnClick speaks it."),
    ("Holding passwords went radioactive",
     "After Amazon v. Perplexity (March 2026), brokering user credentials in the cloud became "
     "a legal minefield. UnClick's hold-nothing design is suddenly the safe one."),
    ("The world moved to signed mandates",
     "Cloudflare, Visa, Mastercard and Google all shifted to user-issued, signed agent "
     "permissions. UnClick is already built that exact shape."),
]
cw2, ch2 = 5.55, 1.78
for i,(t,b) in enumerate(cards):
    cx = ML + (i%2)*(cw2+0.4); cyy = 2.05 + (i//2)*(ch2+0.32)
    rect(s, cx, cyy, cw2, ch2, fill=WHITE, line=LINE, line_w=1.0)
    text(s, cx+0.32, cyy+0.24, cw2-0.6, 0.4, [(t, TEAL_DEEP)], 16, bold=True)
    paras(s, cx+0.32, cyy+0.72, cw2-0.6, ch2-0.85, [b], 12.5, color=BODY, spacing=1.22)
text(s, ML, 6.15, 11.4, 0.5, [("We ride the standard at the precise moment it consolidates.", INK)], 16, bold=True)
footer(s, 5)

# ================================================================ SLIDE 6 helicopter
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "The solution  ·  helicopter view")
text(s, ML, 1.0, 11.5, 1.0, [("UnClick is the ", INK), ("rails", TEAL_DEEP), (" your agent plugs into.", INK)], 38, bold=True, align=PP_ALIGN.CENTER)
text(s, ML, 1.85, 11.5, 0.5, [("One npm install. Five superpowers. Works behind every AI app you already use.", BODY)], 15.5, align=PP_ALIGN.CENTER)
rect(s, 4.97, 2.5, 3.4, 0.62, fill=NAVY, rounded=True, radius=0.25)
text(s, 4.97, 2.5, 3.4, 0.62, [("◆  UnClick · one install", WHITE)], 15, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
pills = [("Apps","Every tool, ready to call."),
         ("Memory","Remembers across every session and device."),
         ("Connections","Sign in once. Every agent inherits."),
         ("AutoPilot","Plans, routes, builds and checks the work."),
         ("XPass","Proves the work before it ships.")]
pw = 2.12; gap = 0.18; total = 5*pw + 4*gap; startx = (W-total)/2; py = 3.5; ph = 2.3
for i,(t,b) in enumerate(pills):
    cx = startx + i*(pw+gap)
    rect(s, cx, py, pw, ph, fill=WHITE, line=LINE, line_w=1.0)
    rect(s, cx, py, pw, 0.06, fill=TEAL, rounded=False)
    rect(s, cx+pw/2-0.28, py+0.28, 0.56, 0.56, fill=RGBColor(0xE4,0xF5,0xF5), line=RGBColor(0xBF,0xE8,0xE9), line_w=1.0, radius=0.25)
    text(s, cx, py+0.34, pw, 0.44, [("◧", TEAL_DEEP)], 18, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx, py+1.0, pw, 0.4, [(t, INK)], 15, bold=True, align=PP_ALIGN.CENTER)
    paras(s, cx+0.16, py+1.42, pw-0.32, 0.8, [b], 11, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.15)
text(s, ML, 6.15, 11.5, 0.5,
     [("Sits behind Claude, ChatGPT, Cursor, Copilot and every MCP client. "
       "We do not compete with them. We make all of them better.", MUTED)], 12.5, align=PP_ALIGN.CENTER)
footer(s, 6)

# ================================================================ SLIDE 7 apps
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Pillar 1  ·  Apps")
text(s, ML, 1.05, 11.4, 1.0, [("An ", INK), ("app store", TEAL_DEEP), (" for AI agents.", INK)], 42, bold=True)
paras(s, ML, 2.0, 10.9, 1.0,
      ["One install, and your agent can actually do things: send an email, charge a card on "
       "Stripe, post to Slack, open a GitHub pull request, check the weather, read the news. "
       "No more wiring up six servers by hand."],
      16, color=BODY, spacing=1.3)
stats = [("450+","Callable endpoints"),("178+","Tools"),("60+","Integrations"),("20+","One-click connections")]
sw = 2.72; sgap = 0.18; sx0 = ML; sy = 3.5; sh = 1.7
for i,(n,l) in enumerate(stats):
    cx = sx0 + i*(sw+sgap)
    rect(s, cx, sy, sw, sh, fill=WHITE, line=LINE, line_w=1.0)
    text(s, cx, sy+0.32, sw, 0.8, [(n, TEAL_DEEP)], 40, bold=True, align=PP_ALIGN.CENTER, name=MONO)
    text(s, cx, sy+1.15, sw, 0.4, [(l, MUTED)], 12, bold=True, align=PP_ALIGN.CENTER)
paras(s, ML, 5.6, 11.4, 0.9,
      ["Your agent discovers and calls any tool on its own. New tools ship to the cloud and "
       "every connected agent picks them up automatically. No package update needed, ever."],
      12.5, color=MUTED, spacing=1.25)
footer(s, 7)

# ================================================================ SLIDE 8 memory
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Pillar 2  ·  Memory")
text(s, ML, 1.0, 11.4, 1.0, [("Your agent forgets everything. ", INK), ("We fix that.", TEAL_DEEP)], 38, bold=True)
layers = [("1","Identity & standing rules","Who you are and your always-on rules, loaded first.", RGBColor(0x1C,0x5A,0x6B)),
          ("2","Knowledge library","Your documents and briefs, neatly versioned.", RGBColor(0x23,0x6D,0x7E)),
          ("3","Session summaries","Picks up exactly where you left off.", RGBColor(0x2A,0x80,0x90)),
          ("4","Searchable facts","Conversations distilled into facts overnight.", RGBColor(0x33,0x91,0x9F)),
          ("5","Full conversation log","Every word, verbatim and searchable.", RGBColor(0x3D,0xA2,0xAC)),
          ("6","Code store","Code kept aside, expanded only on demand.", RGBColor(0x49,0xB3,0xBB))]
lx, lw, ly, lh, lgap = ML, 6.7, 2.1, 0.66, 0.12
for i,(n,t,d,c) in enumerate(layers):
    yy = ly + i*(lh+lgap)
    rect(s, lx, yy, lw, lh, fill=WHITE, line=LINE, line_w=1.0)
    rect(s, lx+0.16, yy+0.16, 0.34, 0.34, fill=c, rounded=True, radius=0.25)
    text(s, lx+0.16, yy+0.13, 0.34, 0.36, [(n, WHITE)], 12, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, name=MONO)
    text(s, lx+0.66, yy+0.07, lw-0.8, 0.32, [(t, INK)], 13.5, bold=True)
    text(s, lx+0.66, yy+0.36, lw-0.8, 0.28, [(d, MUTED)], 10.5)
# side dark card
sx, sw2 = 8.0, 4.42
rect(s, sx, 2.1, sw2, 4.66, fill=NAVY, rounded=True, radius=0.06)
text(s, sx+0.4, 2.5, sw2-0.8, 0.6, [("Your data. Your database.", WHITE)], 22, bold=True)
paras(s, sx+0.4, 3.25, sw2-0.8, 2.2,
      ["UnClick memory lives in your own store, not ours. We never see it. Others keep your "
       "memories in their cloud. We keep them in yours. If you ever leave, nothing moves. "
       "It is already yours."],
      14.5, color=DARKBODY, spacing=1.32)
pill(s, sx+0.4, 5.9, 3.6, 0.5, "Survives new models, new apps, new devices", dark=True)
footer(s, 8)

# ================================================================ SLIDE 9 connections
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Pillar 3  ·  Connections")
text(s, ML, 1.05, 11.4, 1.0, [("Sign in once. ", INK), ("Every agent inherits.", TEAL_DEEP)], 40, bold=True)
paras(s, ML, 2.0, 10.9, 1.0,
      ["Connect your accounts a single time. From then on, any agent on any device can use "
       "them, with no re-logins and no copy-pasting secret keys."],
      16, color=BODY, spacing=1.3)
conn = [("Connect once, works forever","One connection flows to every agent and every device you use."),
        ("We hold mandates, not passwords","Permission to act, not your raw keys. There is almost nothing to steal."),
        ("The shape regulators now require","Built around signed user mandates, the model the industry just standardised on.")]
cw3 = 3.65; cgap = 0.27; cy3 = 3.35; ch3 = 2.4
for i,(t,b) in enumerate(conn):
    cx = ML + i*(cw3+cgap)
    rect(s, cx, cy3, cw3, ch3, fill=WHITE, line=LINE, line_w=1.0)
    rect(s, cx+0.32, cy3+0.3, 0.6, 0.6, fill=RGBColor(0xE4,0xF5,0xF5), line=RGBColor(0xBF,0xE8,0xE9), line_w=1.0, radius=0.22)
    text(s, cx+0.32, cy3+0.32, 0.6, 0.56, [("✓", TEAL_DEEP)], 18, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+0.32, cy3+1.05, cw3-0.6, 0.7, [(t, INK)], 15, bold=True)
    paras(s, cx+0.32, cy3+1.62, cw3-0.6, 0.7, [b], 12, color=BODY, spacing=1.2)
text(s, ML, 6.2, 11.4, 0.5, [("Less for an attacker to steal. Less for you to manage. Safer by design.", INK)], 16, bold=True)
footer(s, 9)

# ================================================================ SLIDE 10 autopilot
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Pillar 4  ·  AutoPilot")
text(s, ML, 1.0, 11.4, 1.0, [("A self-driving ", INK), ("factory", TEAL_DEEP), (" for getting work done.", INK)], 36, bold=True)
paras(s, ML, 1.9, 11.0, 0.8,
      ["Type a request in plain English. It moves down an assembly line of specialist AI workers. "
       "One clear job goes in. One safe result comes out."],
      15.5, color=BODY, spacing=1.28)
# launch bar
rect(s, ML, 2.95, 11.5, 0.7, fill=NAVY, rounded=True, radius=0.12)
text(s, ML+0.3, 2.95, 8, 0.7, [("◐  Launchpad · a human watches every step", WHITE)], 14, bold=True, anchor=MSO_ANCHOR.MIDDLE)
rect(s, ML+9.35, 3.08, 1.95, 0.44, fill=RGBColor(0x3A,0x16,0x16), line=RGBColor(0x8A,0x33,0x2E), line_w=1.2, radius=0.3)
text(s, ML+9.35, 3.08, 1.95, 0.44, [("● KILL SWITCH", RGBColor(0xFF,0x8D,0x83))], 11, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
stations = [("01","Intake","Plain-English request",GREEN),("02","Research","Explore options",GREEN),
            ("03","Plan","A small work packet",GREEN),("04","Build","Focused changes",GREEN),
            ("05","Test","Run the checks",AMBER),("06","Review","A second set of eyes",GREEN),
            ("07","Safety","Guard the danger",GREEN),("08","Ship","Merge, publish, prove",GREEN)]
stw = 1.36; stgap = 0.085; stx0 = ML; sty = 4.0; sth = 1.55
for i,(n,t,d,c) in enumerate(stations):
    cx = stx0 + i*(stw+stgap)
    rect(s, cx, sty, stw, sth, fill=WHITE, line=LINE, line_w=1.0)
    rect(s, cx+stw-0.26, sty+0.16, 0.12, 0.12, fill=c, rounded=True, radius=0.5)
    text(s, cx+0.14, sty+0.18, stw-0.3, 0.3, [(n, TEAL_DEEP)], 10.5, bold=True, name=MONO)
    text(s, cx+0.14, sty+0.5, stw-0.28, 0.32, [(t, INK)], 13, bold=True)
    paras(s, cx+0.14, sty+0.85, stw-0.28, 0.6, [d], 9, color=MUTED, spacing=1.1)
# legend
lx = ML
for lab,c in [("Passed",GREEN),("Waiting",AMBER),("Blocked",RED)]:
    rect(s, lx, 5.95, 0.16, 0.16, fill=c, rounded=True, radius=0.5)
    text(s, lx+0.24, 5.9, 1.3, 0.3, [(lab, BODY)], 12, bold=True)
    lx += 1.35
text(s, lx+0.2, 5.88, 8, 0.3, [("Autonomy with Observability. Command Control at Every Step.", TEAL_DEEP)], 13, bold=True)
footer(s, 10)

# ================================================================ SLIDE 11 safety (dark)
s = slide(); fullbg(s, NAVY)
eyebrow(s, "Trust  ·  safety  ·  proof", y=0.6, dark=True)
text(s, ML, 1.0, 11.5, 1.4, [("Autonomy without observability is a ", WHITE), ("blocker, not a feature.", TEAL_BR)], 34, bold=True, spacing=1.04)
# gate before
gw, gy, gh = 5.0, 2.55, 1.9
rect(s, ML, gy, gw, gh, fill=RGBColor(0x2A,0x25,0x12), line=RGBColor(0x6A,0x53,0x20), line_w=1.2)
text(s, ML+0.35, gy+0.28, gw-0.7, 0.3, [("BEFORE THE ACTION", AMBER)], 11, bold=True)
text(s, ML+0.35, gy+0.62, gw-0.7, 0.4, [("XGate", WHITE)], 22, bold=True)
paras(s, ML+0.35, gy+1.12, gw-0.7, 0.75,
      ["Can this even happen? Gates block secret leaks, destructive database wipes, dangerous "
       "git rewrites, surprise deploys and overspend, before a single line runs."],
      11.5, color=DARKBODY, spacing=1.2)
# action divider
rect(s, ML+gw+0.18, gy+0.4, 0.6, 1.1, fill=RGBColor(0x12,0x36,0x47), line=RGBColor(0x2A,0x55,0x66), line_w=1.0, radius=0.2)
text(s, ML+gw+0.18, gy+0.4, 0.6, 1.1, [("ACTION", WHITE)], 11, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# gate after
gx2 = ML+gw+0.96
rect(s, gx2, gy, gw, gh, fill=RGBColor(0x10,0x33,0x3A), line=TEAL, line_w=1.2)
text(s, gx2+0.35, gy+0.28, gw-0.7, 0.3, [("AFTER THE ACTION", TEAL_BR)], 11, bold=True)
text(s, gx2+0.35, gy+0.62, gw-0.7, 0.4, [("XPass", WHITE)], 22, bold=True)
paras(s, gx2+0.35, gy+1.12, gw-0.7, 0.75,
      ["Was it done right? Hundreds of checks across tests, security, UX, copy, legal and SEO. "
       "It loops until the work is green, then signs off."],
      11.5, color=DARKBODY, spacing=1.2)
# trust row
trust = [("Kill switch","Stop everything, instantly. If the state is unclear, it picks the safest option and halts."),
         ("Flight recorder","A tamper-evident log of who did what, why, and the proof. Replay or roll back anything."),
         ("Separation of duties","The worker that builds never approves its own work. Risky jobs need three roles."),
         ("Human in the loop","Money, secrets, domains and launches always wait for a human yes.")]
tw = 2.72; tgap = 0.18; ty = 4.75; th = 1.45
for i,(t,b) in enumerate(trust):
    cx = ML + i*(tw+tgap)
    rect(s, cx, ty, tw, th, fill=CARD_DARK, line=RGBColor(0x1F,0x4A,0x59), line_w=1.0)
    text(s, cx+0.24, ty+0.2, tw-0.48, 0.4, [(t, WHITE)], 13, bold=True)
    paras(s, cx+0.24, ty+0.6, tw-0.48, 0.8, [b], 10, color=RGBColor(0xA7,0xC9,0xCF), spacing=1.15)
text(s, ML, 6.45, 11.5, 0.5, [("Transparent.  Trusted.  Traceable.", TEAL_BR)], 22, bold=True, align=PP_ALIGN.CENTER)
footer(s, 11, dark=True)

# ================================================================ SLIDE 12 moat bundle (dark)
s = slide(); fullbg(s, NAVY)
rect(s, 7.5, -1.5, 8.5, 7.5, fill=NAVY2, rounded=True, radius=0.5)
eyebrow(s, "The moat  ·  why we win", y=0.85, dark=True)
text(s, ML, 1.25, 11.5, 1.4, [("Everyone has a piece. We have the ", WHITE), ("whole stack.", TEAL_BR)], 38, bold=True, align=PP_ALIGN.CENTER, spacing=1.02)
paras(s, ML+0.6, 2.55, 10.3, 0.9,
      ["Memory startups do memory. Most servers do one integration. Coding tools do code. "
       "UnClick is the only place all five live behind a single install and make each other stronger."],
      15, color=DARKBODY, align=PP_ALIGN.CENTER, spacing=1.28)
bundle = [("Catalog","act, not click"),("Memory","you own it"),("Connections","sign in once"),
          ("AutoPilot","moves work"),("Proof","done right")]
bw = 1.78; bgap = 0.52; total = 5*bw + 4*bgap + bw + bgap; startx = (W-total)/2; by = 3.85; bh = 1.25
x = startx
for i,(t,d) in enumerate(bundle):
    rect(s, x, by, bw, bh, fill=CARD_DARK, line=RGBColor(0x2A,0x66,0x70), line_w=1.0)
    text(s, x, by+0.28, bw, 0.4, [(t, WHITE)], 14, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, by+0.72, bw, 0.3, [(d, RGBColor(0x9F,0xC1,0xC7))], 10, align=PP_ALIGN.CENTER)
    op = "+" if i < 4 else "="
    text(s, x+bw, by, bgap, bh, [(op, TEAL_BR)], 22, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += bw + bgap
rect(s, x, by, bw, bh, fill=TEAL, rounded=True, radius=0.12)
text(s, x, by+0.24, bw, 0.4, [("One install", NAVY)], 14, bold=True, align=PP_ALIGN.CENTER)
text(s, x, by+0.7, bw, 0.3, [("the whole stack", NAVY)], 10, align=PP_ALIGN.CENTER)
text(s, ML, 5.6, 11.5, 0.6, [("Stripe for agents, not Windows for agents. We are the rails every harness runs on.", WHITE)], 18, bold=True, align=PP_ALIGN.CENTER)
footer(s, 12, dark=True)

# ================================================================ SLIDE 13 vs table
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Why UnClick vs going it alone")
text(s, ML, 1.0, 11.5, 1.2, [("Claude or ChatGPT can ", INK), ("think.", TEAL_DEEP), (" They cannot remember, connect, or prove.", INK)], 30, bold=True, spacing=1.05)
rows = [("What an agent needs","Claude / ChatGPT / Cursor on their own","With UnClick"),
        ("Memory across sessions & tools","Resets each chat","Remembers everywhere"),
        ("Real tools (act, not click)","Limited or none","450+ endpoints"),
        ("Sign in once across accounts","Re-auth every time","Connect once"),
        ("Works across every AI app","Locked to one vendor","Any MCP client"),
        ("A team of agents that coordinate","One model, alone","Governed worker fleet"),
        ("Safety gates & an audit trail","Trust the prompt","Gates + flight recorder"),
        ("Your data stays yours","Their cloud","Your database")]
tx, ty, tw, th = ML, 2.25, 11.5, 4.05
tbl = s.shapes.add_table(len(rows), 3, Inches(tx), Inches(ty), Inches(tw), Inches(th)).table
tbl.columns[0].width = Inches(4.3); tbl.columns[1].width = Inches(4.1); tbl.columns[2].width = Inches(3.1)
for r,(a,b,c) in enumerate(rows):
    for cidx,val in enumerate((a,b,c)):
        cell = tbl.cell(r,cidx)
        cell.margin_left = Inches(0.16); cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfc = cell.text_frame; tfc.word_wrap = True
        p = tfc.paragraphs[0]; run = p.add_run(); run.text = val
        run.font.name = FONT; run.font.size = Pt(13 if r else 12)
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF1,0xF8,0xF9) if cidx<2 else RGBColor(0xDF,0xF1,0xF1)
            run.font.bold = True; run.font.color.rgb = INK if cidx<2 else TEAL_DEEP
            run.font.size = Pt(12)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if cidx<2 else RGBColor(0xF3,0xFB,0xFB)
            if cidx == 0: run.font.bold = True; run.font.color.rgb = INK
            elif cidx == 1: run.font.bold = True; run.font.color.rgb = RED
            else: run.font.bold = True; run.font.color.rgb = TEAL_DEEP
text(s, ML, 6.5, 11.5, 0.4, [("UnClick is not another chatbot. It is the layer that makes the chatbot you already pay for actually get things done.", MUTED)], 12)
footer(s, 13)

# ================================================================ SLIDE 14 why labs cant
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Defensibility")
text(s, ML, 1.05, 11.4, 1.0, [("The one position the AI labs ", INK), ("structurally cannot take.", TEAL_DEEP)], 34, bold=True)
paras(s, ML, 2.0, 11.0, 1.0,
      ["UnClick's customer is the user, not any single model. An independent layer that audits "
       "and holds every model accountable, including Anthropic's and OpenAI's, is a trust feature "
       "for users and a conflict of interest for the labs."],
      15.5, color=BODY, spacing=1.3)
labs = [("Neutral by design","We back every harness, so we are welcome inside all of them. We do not pick a side in the model wars."),
        ("A compounding data moat","Our edge grows on data only we hold: our own run outcomes and proof history across every model."),
        ("We win on trust, not size","Integration count is a commodity. We compete on curation, reliability and proof that work was done right.")]
cw4 = 3.65; cgap = 0.27; cy4 = 3.4; ch4 = 2.3
for i,(t,b) in enumerate(labs):
    cx = ML + i*(cw4+cgap)
    rect(s, cx, cy4, cw4, ch4, fill=WHITE, line=LINE, line_w=1.0)
    text(s, cx+0.3, cy4+0.28, cw4-0.6, 0.6, [(t, TEAL_DEEP)], 16, bold=True)
    paras(s, cx+0.3, cy4+1.0, cw4-0.6, 1.1, [b], 12, color=BODY, spacing=1.22)
text(s, ML, 6.2, 11.4, 0.5, [("The place where any frontier model is held accountable.", INK)], 16, bold=True)
footer(s, 14)

# ================================================================ SLIDE 15 business model
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Business model")
text(s, ML, 1.0, 11.4, 1.2, [("Simple subscriptions, plus a marketplace that compounds.", INK)], 34, bold=True, spacing=1.05)
rect(s, ML, 2.25, 0.06, 0.5, fill=TEAL, rounded=False)
text(s, ML+0.22, 2.22, 11, 0.5, [("You pay OpenAI for the AI. You pay UnClick for the agent infrastructure.", INK)], 17, bold=True)
tiers = [("FREE","$0","","The full toolbox. 100 calls a day. Self-hosted memory.", False),
         ("PRO","$29","/mo","Unlimited calls. Managed memory with nightly fact extraction and a dashboard.", True),
         ("TEAM","$79","/mo","Up to 5 seats. Shared business context and team observability.", False),
         ("ENTERPRISE","Custom","","SSO, on-prem and custom tooling. Sales-led.", False)]
tw5 = 2.72; tgap5 = 0.18; ty5 = 3.2; th5 = 2.45
for i,(nm,pr,suf,desc,pop) in enumerate(tiers):
    cx = ML + i*(tw5+tgap5)
    rect(s, cx, ty5, tw5, th5, fill=WHITE, line=(TEAL if pop else LINE), line_w=(2.0 if pop else 1.0))
    if pop:
        rect(s, cx+0.3, ty5-0.16, 1.5, 0.32, fill=TEAL, rounded=True, radius=0.5)
        text(s, cx+0.3, ty5-0.17, 1.5, 0.32, [("Most popular", NAVY)], 10, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+0.28, ty5+0.3, tw5-0.5, 0.3, [(nm, TEAL_DEEP)], 12, bold=True)
    text(s, cx+0.28, ty5+0.62, tw5-0.5, 0.6, [(pr, INK),(suf, MUTED)], 30, bold=True, name=MONO)
    paras(s, cx+0.28, ty5+1.42, tw5-0.5, 0.95, [desc], 11.5, color=MUTED, spacing=1.2)
paras(s, ML, 5.95, 11.4, 0.9,
      [[("The second engine: ", TEAL_DEEP, True),
        ("a developer marketplace. Builders ship tools and keep 80% of what they earn; UnClick "
         "takes 20%. Subscription-only billing means we never resell AI tokens, so provider price "
         "swings cannot touch our margin.", MUTED)]],
      12.5, spacing=1.25)
footer(s, 15)

# ================================================================ SLIDE 16 market
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Market")
text(s, ML, 1.05, 11.4, 1.0, [("Every agent user. ", INK), ("Whichever AI they use.", TEAL_DEEP)], 38, bold=True)
aud = [("Businesses & creators","Millions already pay for Claude, ChatGPT and Cursor and do not want to wire up six servers by hand. They want their agent to just work."),
       ("Developers","Builders who want distribution and a way to earn from the tools they make. UnClick is the shelf and the checkout, all at once.")]
aw = 5.55; agap = 0.4; ay = 2.35; ah = 2.4
for i,(t,b) in enumerate(aud):
    cx = ML + i*(aw+agap)
    rect(s, cx, ay, aw, ah, fill=WHITE, line=LINE, line_w=1.0)
    text(s, cx+0.35, ay+0.3, aw-0.7, 0.5, [(t, INK)], 21, bold=True)
    paras(s, cx+0.35, ay+1.0, aw-0.7, 1.2, [b], 14.5, color=BODY, spacing=1.28)
rect(s, ML, 5.1, 11.5, 1.35, fill=RGBColor(0xE9,0xF6,0xF6), line=RGBColor(0xB9,0xE6,0xE7), line_w=1.0)
paras(s, ML+0.4, 5.35, 10.7, 0.9,
      [[("Because we are the ", INK, False),("substrate, not a competitor", INK, True),
        (", our market is the entire agent ecosystem, and it grows every single time a new "
         "harness or a new model launches.", INK, False)]],
      16, spacing=1.3)
footer(s, 16)

# ================================================================ SLIDE 17 traction
s = slide(); fullbg(s, LIGHTBG); accent_top(s)
eyebrow(s, "Where we are today")
text(s, ML, 1.05, 11.4, 1.0, [("Live, in the wild, at ", INK), ("unclick.world", TEAL_DEEP)], 38, bold=True)
live = [("The catalog is live","450+ endpoints across 60+ integrations, callable today."),
        ("Memory engine in production","Six-layer, user-owned, cross-session memory."),
        ("Works across every major client","Claude, ChatGPT, Cursor, Copilot, Windsurf and more."),
        ("AutoPilot fleet running","Governed multi-agent pipeline with safety gates and audit."),
        ("Pricing live","Free, Pro, Team and Enterprise tiers shipped."),
        ("Australian made","Built in Melbourne, with data-sovereignty options.")]
lw6 = 5.55; lgap6 = 0.4; ly6 = 2.2; lh6 = 0.95
for i,(t,b) in enumerate(live):
    cx = ML + (i%2)*(lw6+lgap6); yy = ly6 + (i//2)*(lh6+0.18)
    rect(s, cx, yy, 0.42, 0.42, fill=RGBColor(0xE3,0xF6,0xEC), line=RGBColor(0xA9,0xE2,0xC4), line_w=1.0, radius=0.22)
    text(s, cx, yy, 0.42, 0.42, [("✓", GREEN)], 14, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+0.58, yy-0.02, lw6-0.7, 0.34, [(t, INK)], 15, bold=True)
    text(s, cx+0.58, yy+0.34, lw6-0.7, 0.5, [(b, MUTED)], 12)
rect(s, ML, 5.95, 11.5, 0.62, fill=RGBColor(0xE9,0xF6,0xF6), line=TEAL, line_w=1.2)
text(s, ML, 5.95, 11.5, 0.62,
     [("Add your live traction here  →  active users · calls per month · MRR · week-on-week growth", TEAL_DEEP)],
     13.5, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 17)

# ================================================================ SLIDE 18 vision (dark)
s = slide(); fullbg(s, NAVY)
rect(s, 7.0, -1.5, 9.0, 7.5, fill=NAVY2, rounded=True, radius=0.5)
eyebrow(s, "The vision", y=0.95, dark=True)
text(s, ML, 1.4, 11.3, 1.4, [("The trust layer for the entire ", WHITE), ("agent economy.", TEAL_BR)], 40, bold=True, spacing=1.04)
paras(s, ML, 3.1, 10.5, 1.3,
      ["Today, UnClick is the invisible plumbing that makes agents useful. Tomorrow, it is the "
       "place where any AI, from any lab, is held accountable: every action proven, every memory "
       "owned by you, every job traceable."],
      16.5, color=DARKBODY, spacing=1.34)
text(s, ML, 4.95, 11.3, 1.2, [("Where AI belongs. ", WHITE), ("Humans welcome.", TEAL_BR)], 44, bold=True, spacing=1.02)
footer(s, 18, dark=True)

# ================================================================ SLIDE 19 ask (dark)
s = slide(); fullbg(s, NAVY)
logo_wordmark(s, ML, 0.7, dark=True, size=28)
eyebrow(s, "Join us", y=1.65, dark=True)
text(s, ML, 2.05, 11.3, 1.0, [("Let's build ", WHITE), ("where AI belongs.", TEAL_BR)], 40, bold=True)
asks = [("01","Grow the catalog","more tools, more integrations, more one-click connections."),
        ("02","Open the marketplace","turn on developer payouts and the two-sided flywheel."),
        ("03","Scale the AutoPilot fleet","more governed autonomy, more proof, more throughput.")]
ay = 3.15
for n,t,d in asks:
    text(s, ML, ay, 0.6, 0.4, [(n, TEAL_BR)], 16, bold=True, name=MONO)
    text(s, ML+0.7, ay, 10.4, 0.4, [(t+"  ·  ", WHITE, True),(d, DARKBODY)], 15)
    ay += 0.56
rect(s, ML, 5.1, 11.0, 0.6, fill=RGBColor(0x10,0x33,0x44), line=TEAL, line_w=1.2)
text(s, ML, 5.1, 11.0, 0.6, [("Add your ask here  →  we are raising $X to reach Y over the next N months", TEAL_BR)], 13.5, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, ML, 6.0, 11.3, 0.5,
     [("unclick.world    |    creativelead@malamutemayhem.com    |    Built in Melbourne, Australia", DARKBODY)],
     14, bold=True)
footer(s, 19, dark=True)

# ---------------------------------------------------------------- save
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "unclick-investor-deck.pptx")
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
